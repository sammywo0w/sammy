from flask import Flask, request, send_file, jsonify
from pptx import Presentation
import os
import uuid
import re

app = Flask(__name__)

@app.route('/generate', methods=['POST'])
def generate_pptx():
    data = request.get_json()
    if not data:
        return jsonify({"error": "Invalid or missing JSON"}), 400

    prs = Presentation("11Business Model Template.pptx")

    placeholder_map = {
        "target_customers": "Target Customer",
        "value_proposition": "Value Proposition",
        "cost_drivers": "Cost driver",
        "revenue_drivers": "Revenue driver",
        "customer_channels": "Customer Channel",
        "resources": "Resource",
        "offering": "Offering",
        "customer_relationship": "Customer Relation",
        "partners": "Partner",
        "activities": "Activity",
        "model_name": "Business Model Name"
    }

    def replace_placeholders_by_index(slide, key, items):
        placeholder_base = placeholder_map[key]
        max_range = 30  # защита от слишком больших циклов

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    original = run.text
                    if not isinstance(original, str):
                        continue

                    # Заменим пронумерованные <Key 1>, <Key 2> ...
                    for i in range(1, max_range):
                        placeholder = f"<{placeholder_base} {i}>"
                        if placeholder in run.text:
                            new_value = items[i - 1] if i <= len(items) else ""
                            run.text = run.text.replace(placeholder, new_value)

                    # Также заменим одиночные плейсхолдеры, типа <Business Model Name>
                    if f"<{placeholder_base}>" in run.text:
                        text_value = items[0] if items else ""
                        run.text = run.text.replace(f"<{placeholder_base}>", text_value)

    # Применяем для всех слайдов и ключей
    for slide in prs.slides:
        for key in placeholder_map:
            values = data.get(key, [])
            if not isinstance(values, list):
                values = [values]
            replace_placeholders_by_index(slide, key, values)

    # Сохраняем итог
    filename = f"{uuid.uuid4().hex}.pptx"
    filepath = os.path.join("/tmp", filename)
    prs.save(filepath)

    return send_file(filepath, as_attachment=True, download_name="business_model_canvas.pptx")

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=int(os.environ.get("PORT", 5000)))
